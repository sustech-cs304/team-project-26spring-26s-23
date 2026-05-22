"""File-to-text conversion helpers for runtime tools."""

from __future__ import annotations

from pathlib import Path

import pdfplumber
from docx import Document
from pptx import Presentation


def convert_file_to_str(
    path: str | Path,
    *,
    max_chars: int | None = None,
    encoding: str = "utf-8",
) -> str:
    resolved = Path(path).expanduser()
    if resolved.is_absolute():
        workspace_root = Path.cwd().resolve(strict=False)
        target = resolved.resolve(strict=False)
        try:
            target.relative_to(workspace_root)
        except ValueError as exc:
            raise ValueError("path must be within the workspace.") from exc
    else:
        workspace_root = Path.cwd().resolve(strict=False)
        target = (workspace_root / resolved).resolve(strict=False)
    if not target.exists():
        raise FileNotFoundError(str(target))
    if not target.is_file():
        raise ValueError("path must be a file.")

    suffix = target.suffix.lower()
    if suffix == ".pdf":
        text = _convert_pdf_to_text(target)
    elif suffix == ".docx":
        text = _convert_docx_to_text(target)
    elif suffix == ".pptx":
        text = _convert_pptx_to_text(target)
    else:
        text = target.read_text(encoding=encoding, errors="replace")

    normalized = text.strip()
    if max_chars is not None and max_chars > 0 and len(normalized) > max_chars:
        return normalized[:max_chars]
    return normalized


def _convert_pdf_to_text(path: Path) -> str:
    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        for index, page in enumerate(pdf.pages, start=1):
            extracted = (page.extract_text() or "").strip()
            if extracted:
                parts.append(f"[Page {index}]\n{extracted}")
    return "\n\n".join(parts)


def _convert_docx_to_text(path: Path) -> str:
    doc = Document(str(path))
    parts = [paragraph.text.strip() for paragraph in doc.paragraphs]
    return "\n".join(item for item in parts if item)


def _convert_pptx_to_text(path: Path) -> str:
    presentation = Presentation(str(path))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text_frame = getattr(shape, "text_frame", None)
            if text_frame is None:
                continue
            text = (shape.text or "").strip()
            if text:
                parts.append(text)
    return "\n".join(parts)

